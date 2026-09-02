"""Build-time gate: every accepted upload format must survive this image's libmagic.

## Why this exists

``file`` 5.46 (Debian trixie ``1:5.46-5``, the version ``python:3.13-slim``
resolved to before the pin) regressed on the BUFFER entry point: a plain
``PK\\x03\\x04`` archive came back as ``application/octet-stream`` / ``data``
from ``magic_buffer`` while ``magic_file`` on the very same bytes returned
``application/zip`` (Debian #1102577; the cherry-pick in ``-5`` does not cover
it). Stage 1 reads uploads through ``magic.from_buffer``, so every student
archive was rejected as ``magic_mismatch`` in production while every local
check stayed green -- the host's libmagic was a different build.

The lesson is not "pin 5.47". It is that a silent detector regression is
invisible to every test that does not run **inside the built image**, and that
the blast radius of a libmagic bump is the whole upload surface, not just the
format that happened to break. So: this gate runs at image build time, against
the image's own libmagic, over every format the system accepts.

## What it checks

Coverage is DERIVED, never hand-listed:

    _EXTENSION_TO_MIME_FAMILIES.keys() | HOMEWORK_POLICY.allowed_extensions

Adding a format to either set without teaching this script to build a fixture
for it turns the build red. That is deliberate -- it guards the construction,
the same way ``test_code_extensions_are_text_in_normalizer`` guards the derived
union in ``normalizer/classify.py`` rather than its values.

Each format is checked through :func:`verify_extension_matches_content` -- the
exact function Stage 1 calls, not a re-implementation. A gate that probed
``magic.from_buffer`` directly could stay green while the real predicate broke.

## Why fixtures are generated, not committed

They must be real files (impl-rules#13): a hand-written byte stub proves
nothing about a detector that reads container structure. Generating them from
the libraries already in the image (stdlib archives, Pillow, PyMuPDF,
python-docx, python-pptx) and ffmpeg keeps them real, deterministic, free of
binary blobs in git, and -- because the whole run happens in one build layer
that ends in ``rm -rf`` -- absent from the shipped image.

The legacy CDFV2 containers are the exception, and not uniformly: measured in
the built images 2026-09-02, ``libreoffice-impress`` (the only LibreOffice in
either image, and only in ``Dockerfile.worker``) converts ``pptx -> ppt`` but
NOT ``docx -> doc`` or ``xlsx -> xls`` -- those need Writer and Calc, which are
installed nowhere. What cannot be proven is therefore declared per image in
:data:`UNPROVABLE`, one table with a reason per entry rather than conditionals
scattered through the run: a format the gate passes over must always be able to
say why, in the build log, at the moment it is passed over.

Usage::

    python scripts/magic_format_gate.py --image worker   # Dockerfile.worker
    python scripts/magic_format_gate.py --image app      # Dockerfile
"""

from __future__ import annotations

import argparse
import functools
import io
import shutil
import subprocess
import sys
import tarfile
import tempfile
import zipfile
from collections.abc import Callable
from pathlib import Path

from course_supporter.security.exceptions import SecurityRejectedError
from course_supporter.security.file_type import (
    _EXTENSION_TO_MIME_FAMILIES,
    detect_mime_type,
    verify_extension_matches_content,
)
from course_supporter.security.policies import HOMEWORK_POLICY

# Formats this gate cannot prove, declared per image with the reason it cannot.
# Everything NOT listed here must pass; an unlisted format that fails fails the
# build. Keeping the exceptions in one table (rather than as conditionals inside
# the run) is what makes them auditable: adding a skip means writing down why.
#
# The reason strings are operator-facing and print into the build log verbatim.
_NO_WRITER_CALC = (
    "no Writer/Calc in either image; no policy (AUTHORED, HOMEWORK) accepts it "
    "— unreachable for upload"
)
_NO_LIBREOFFICE = "no LibreOffice in the app image; proven in the worker image"

UNPROVABLE: dict[str, dict[str, str]] = {
    # Dockerfile — no LibreOffice at all, so every legacy CDFV2 target is out.
    "app": {
        "doc": _NO_WRITER_CALC,
        "xls": _NO_WRITER_CALC,
        "ppt": _NO_LIBREOFFICE,
    },
    # Dockerfile.worker — libreoffice-impress proves ppt; Writer and Calc are
    # absent, so doc and xls stay out here too.
    "worker": {
        "doc": _NO_WRITER_CALC,
        "xls": _NO_WRITER_CALC,
    },
}

# Legacy CDFV2 targets that a LibreOffice module can produce when present.
LEGACY_OFFICE: frozenset[str] = frozenset({"doc", "ppt", "xls"})

_TEXT_SAMPLE = (
    "# Заголовок роботи\n\nПерший абзац українською.\nSecond line in English.\n"
)

# Per-language snippets for the code extensions. Each is real source, not a
# placeholder: the textual invariant probes the charset, and a snippet that
# looks like the language keeps the fixture honest if the invariant is ever
# tightened to sniff syntax.
_CODE_SNIPPETS: dict[str, str] = {
    "py": "def greet(name: str) -> str:\n    return f'Привіт, {name}'\n",
    "ipynb": '{"cells": [], "metadata": {}, "nbformat": 4, "nbformat_minor": 5}\n',
    "js": "export function greet(name) {\n  return `Hi ${name}`;\n}\n",
    "mjs": "export const greet = (name) => `Hi ${name}`;\n",
    "cjs": "module.exports = function greet(name) {\n  return 'Hi ' + name;\n};\n",
    "jsx": "export const App = () => <main>Привіт</main>;\n",
    "ts": "export function greet(name: string): string {\n  return `Hi ${name}`;\n}\n",
    "tsx": "export const App = (): JSX.Element => <main>Привіт</main>;\n",
    "java": "public class Main {\n  public static void main(String[] a) {}\n}\n",
    "kt": 'fun greet(name: String) = "Привіт, $name"\n',
    "kts": 'val greeting: String = "Привіт"\nprintln(greeting)\n',
    "cs": "public class Program {\n  static void Main() {}\n}\n",
    "go": 'package main\n\nimport "fmt"\n\nfunc main() { fmt.Println("Привіт") }\n',
    "rs": 'fn main() {\n    println!("Привіт");\n}\n',
    "php": "<?php\nfunction greet(string $n): string { return 'Hi ' . $n; }\n",
    "rb": 'def greet(name)\n  "Привіт, #{name}"\nend\n',
    "c": '#include <stdio.h>\n\nint main(void) { puts("Привіт"); return 0; }\n',
    "h": "#ifndef GREET_H\n#define GREET_H\nvoid greet(const char *name);\n#endif\n",
    "cpp": '#include <iostream>\n\nint main() { std::cout << "Привіт\\n"; }\n',
    "hpp": "#pragma once\n#include <string>\nstd::string greet(std::string name);\n",
    "cc": "#include <string>\n\nstd::string greet(std::string n) { return n; }\n",
    "swift": 'func greet(_ name: String) -> String { "Привіт, \\(name)" }\n',
    "dart": "String greet(String name) => 'Привіт, $name';\n",
    "html": "<!doctype html>\n<html lang='uk'><body><p>Привіт</p></body></html>\n",
    "htm": "<!doctype html>\n<html lang='uk'><body><p>Привіт</p></body></html>\n",
    "css": "body {\n  font-family: system-ui;\n  color: #222;\n}\n",
    "scss": "$brand: #222;\n\nbody {\n  color: $brand;\n}\n",
    "json": '{\n  "title": "Робота",\n  "score": 100\n}\n',
    "xml": (
        '<?xml version="1.0" encoding="UTF-8"?>\n<work><title>Робота</title></work>\n'
    ),
    "yaml": "title: Робота\nscore: 100\n",
    "yml": "title: Робота\nscore: 100\n",
    "toml": '[work]\ntitle = "Робота"\nscore = 100\n',
    "sql": "SELECT id, title\nFROM works\nWHERE score > 50;\n",
    "sh": '#!/bin/sh\nset -e\necho "Привіт"\n',
    "md": _TEXT_SAMPLE,
    "markdown": _TEXT_SAMPLE,
    "txt": _TEXT_SAMPLE,
    "csv": "id,title,score\n1,Робота,100\n2,Друга робота,80\n",
}

# ffmpeg recipes: (input spec, candidate encoders tried in order). Durations and
# frame sizes are the smallest that still produce a well-formed container --
# libmagic only reads the head, but the head must be a real one.
#
# Encoders are a CHAIN rather than a single choice because which optional codecs
# ffmpeg carries differs per build (Debian's ships libvorbis; the operator's
# homebrew build does not, and fails with exit 8). The chain keeps one recipe
# honest across both images and the dev machine; exhausting it fails the gate
# loudly rather than skipping the format.
_SINE = ["-f", "lavfi", "-i", "sine=frequency=440:duration=0.3"]
_TESTSRC = ["-f", "lavfi", "-i", "testsrc=duration=0.3:size=64x64:rate=10"]

_FFMPEG_RECIPES: dict[str, tuple[list[str], tuple[list[str], ...]]] = {
    "mp3": (_SINE, (["-c:a", "libmp3lame"], ["-c:a", "mp3"])),
    "wav": (_SINE, (["-c:a", "pcm_s16le"],)),
    "ogg": (_SINE, (["-c:a", "libvorbis"], ["-c:a", "libopus"])),
    "flac": (_SINE, (["-c:a", "flac"],)),
    "m4a": (_SINE, (["-c:a", "aac"],)),
    "mp4": (_TESTSRC, (["-c:v", "mpeg4"],)),
    "mov": (_TESTSRC, (["-c:v", "mpeg4"],)),
    "avi": (_TESTSRC, (["-c:v", "mpeg4"],)),
    "mkv": (_TESTSRC, (["-c:v", "mpeg4"],)),
    "webm": (_TESTSRC, (["-c:v", "libvpx"], ["-c:v", "vp8", "-strict", "-2"])),
}


def _tool(name: str) -> str:
    """Resolve an external tool to a full path, or fail with a named reason."""
    resolved = shutil.which(name)
    if resolved is None:
        raise RuntimeError(f"{name} is not on PATH in this image")
    return resolved


def _write_text(path: Path, *, ext: str) -> None:
    path.write_text(_CODE_SNIPPETS[ext], encoding="utf-8")


def _write_zip(path: Path) -> None:
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("main.py", _CODE_SNIPPETS["py"])
        zf.writestr("README.md", _TEXT_SAMPLE)


def _write_tar(path: Path, *, compress: bool) -> None:
    mode = "w:gz" if compress else "w"
    with tarfile.open(path, mode) as tf:  # type: ignore[call-overload]
        members = (("main.py", _CODE_SNIPPETS["py"]), ("README.md", _TEXT_SAMPLE))
        for name, body in members:
            raw = body.encode("utf-8")
            info = tarfile.TarInfo(name)
            info.size = len(raw)
            tf.addfile(info, io.BytesIO(raw))


def _write_pdf(path: Path) -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 96), "Robota studenta")
    doc.save(str(path))
    doc.close()


def _write_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_paragraph("Текст роботи студента.")
    document.save(str(path))


def _write_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))


def _write_xlsx(path: Path) -> None:
    """Minimal but structurally real OOXML spreadsheet package.

    ``openpyxl`` is not a runtime dependency, and libmagic identifies xlsx by
    the OOXML package shape (``[Content_Types].xml`` plus an ``xl/`` member),
    so assembling that package directly is both sufficient for the detector and
    free of a new dependency.
    """
    ns_pkg = "http://schemas.openxmlformats.org/package/2006"
    ns_doc = "http://schemas.openxmlformats.org/officeDocument/2006"
    ns_sml = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    xml_head = '<?xml version="1.0" encoding="UTF-8"?>'
    ct = (
        f"{xml_head}"
        f'<Types xmlns="{ns_pkg}/content-types">'
        f'<Default Extension="rels" '
        f'ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        f'<Override PartName="/xl/workbook.xml" ContentType="application/vnd.'
        f'openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
        f"</Types>"
    )
    rels = (
        f"{xml_head}"
        f'<Relationships xmlns="{ns_pkg}/relationships">'
        f'<Relationship Id="rId1" Type="{ns_doc}/relationships/officeDocument" '
        f'Target="xl/workbook.xml"/>'
        f"</Relationships>"
    )
    workbook = (
        f"{xml_head}"
        f'<workbook xmlns="{ns_sml}">'
        f'<sheets><sheet name="Sheet1" sheetId="1" r:id="rId1" '
        f'xmlns:r="{ns_doc}/relationships"/></sheets>'
        f"</workbook>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        # ``[Content_Types].xml`` must be the first member of an OOXML package.
        zf.writestr("[Content_Types].xml", ct)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("xl/workbook.xml", workbook)


def _write_image(path: Path, *, fmt: str) -> None:
    from PIL import Image

    Image.new("RGB", (16, 16), (32, 64, 128)).save(path, format=fmt)


def _write_ffmpeg(path: Path, *, ext: str) -> None:
    source, encoder_chain = _FFMPEG_RECIPES[ext]
    ffmpeg = _tool("ffmpeg")
    errors: list[str] = []
    for encoder in encoder_chain:
        result = subprocess.run(
            [
                ffmpeg,
                "-hide_banner",
                "-loglevel",
                "error",
                "-y",
                *source,
                *encoder,
                str(path),
            ],
            capture_output=True,
            timeout=120,
        )
        if result.returncode == 0 and path.exists() and path.stat().st_size > 0:
            return
        errors.append(f"{' '.join(encoder)} -> rc={result.returncode}")
    raise RuntimeError(f"no working encoder for .{ext}: " + "; ".join(errors))


# LibreOffice converts FROM a modern OOXML file, so each legacy target names the
# generated source it is produced from.
_LEGACY_SOURCE: dict[str, tuple[str, Callable[[Path], None]]] = {
    "doc": ("docx", _write_docx),
    "ppt": ("pptx", _write_pptx),
    "xls": ("xlsx", _write_xlsx),
}


def _write_legacy_office(path: Path, *, ext: str, workdir: Path) -> None:
    """Convert a generated OOXML fixture into its legacy CDFV2 counterpart."""
    modern_ext, build_modern = _LEGACY_SOURCE[ext]
    source = workdir / f"legacy-source.{modern_ext}"
    build_modern(source)
    subprocess.run(
        [
            _tool("soffice"),
            "--headless",
            "--convert-to",
            ext,
            "--outdir",
            str(workdir),
            str(source),
        ],
        check=True,
        capture_output=True,
        timeout=180,
    )
    produced = source.with_suffix(f".{ext}")
    produced.replace(path)


def _builders(workdir: Path) -> dict[str, Callable[[Path], None]]:
    """Map every covered extension to the callable that materialises it."""
    part = functools.partial
    builders: dict[str, Callable[[Path], None]] = {
        ext: part(_write_text, ext=ext) for ext in _CODE_SNIPPETS
    }
    builders |= {
        "zip": _write_zip,
        "tar": part(_write_tar, compress=False),
        "gz": part(_write_tar, compress=True),
        "tgz": part(_write_tar, compress=True),
        "pdf": _write_pdf,
        "docx": _write_docx,
        "pptx": _write_pptx,
        "xlsx": _write_xlsx,
        "png": part(_write_image, fmt="PNG"),
        "gif": part(_write_image, fmt="GIF"),
        "webp": part(_write_image, fmt="WEBP"),
        "jpg": part(_write_image, fmt="JPEG"),
        "jpeg": part(_write_image, fmt="JPEG"),
    }
    builders |= {ext: part(_write_ffmpeg, ext=ext) for ext in _FFMPEG_RECIPES}
    builders |= {
        ext: part(_write_legacy_office, ext=ext, workdir=workdir)
        for ext in LEGACY_OFFICE
    }
    return builders


def covered_extensions() -> frozenset[str]:
    """Every extension the gate must prove, derived from the live sources."""
    return frozenset(_EXTENSION_TO_MIME_FAMILIES) | HOMEWORK_POLICY.allowed_extensions


def run_gate(*, image: str, workdir: Path) -> int:
    """Build a fixture per covered extension and push it through Stage 1's check.

    Returns the process exit code: 0 when every format is recognised.
    """
    covered = covered_extensions()
    builders = _builders(workdir)
    unprovable = UNPROVABLE[image]

    missing = sorted(covered - set(builders))
    failures: list[str] = []

    print(
        f"magic gate: {len(covered)} covered extensions, "
        f"checked through verify_extension_matches_content (from_buffer)"
    )
    for ext, reason in sorted(unprovable.items()):
        print(f"  skip {ext:9s} {reason}")

    for ext in sorted(covered):
        if ext in unprovable:
            continue
        if ext not in builders:
            continue  # reported via ``missing`` below
        target = workdir / f"gate-sample.{ext}"
        try:
            builders[ext](target)
        except Exception as exc:  # any fixture failure must fail the gate
            failures.append(f"{ext}: fixture build failed: {type(exc).__name__}: {exc}")
            continue

        data = target.read_bytes()
        if not data:
            failures.append(f"{ext}: fixture is empty")
            continue

        try:
            verify_extension_matches_content(target.name, data)
        except SecurityRejectedError as exc:
            failures.append(f"{ext}: {exc.category.value}: {exc.detail}")
            continue

        print(f"  ok  {ext:9s} {len(data):>8} B  {detect_mime_type(data)}")

    if missing:
        failures.append(
            "no fixture builder for: "
            + ", ".join(missing)
            + " — a format was added to the policy or the family table without"
            " teaching this gate to produce a real sample for it"
        )

    if failures:
        print("\nmagic gate FAILED:", file=sys.stderr)
        for line in failures:
            print(f"  - {line}", file=sys.stderr)
        return 1

    proven = len(covered) - len(unprovable)
    print(f"magic gate OK: {proven} formats recognised by libmagic in this image")
    return 0


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--image",
        choices=tuple(UNPROVABLE),
        required=True,
        help="Which image is being built; selects its UNPROVABLE exception list.",
    )
    args = parser.parse_args()

    workdir = Path(tempfile.mkdtemp(prefix="magic-gate-"))
    try:
        return run_gate(image=args.image, workdir=workdir)
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


if __name__ == "__main__":
    raise SystemExit(main())
