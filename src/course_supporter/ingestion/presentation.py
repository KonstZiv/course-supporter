"""Presentation processor for PDF and PPTX files."""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

import structlog

from course_supporter.ingestion.base import (
    SourceProcessor,
    UnsupportedFormatError,
)
from course_supporter.ingestion.heavy_steps import (
    DescribeSlidesFunc,
    DescribeSlidesParams,
)
from course_supporter.models.source import (
    ChunkType,
    ContentChunk,
    SourceDocument,
    SourceType,
)

if TYPE_CHECKING:
    from course_supporter.llm.router import ModelRouter
    from course_supporter.storage.orm import SourceMaterial

logger = structlog.get_logger()

SUPPORTED_EXTENSIONS = {".pdf", ".pptx"}


class PresentationProcessor(SourceProcessor):
    """Process PDF and PPTX presentations.

    Extracts text from each slide/page as SLIDE_TEXT chunks.
    Optionally uses an injected ``DescribeSlidesFunc`` to describe
    slide images as SLIDE_DESCRIPTION chunks.
    """

    def __init__(
        self,
        *,
        describe_slides_func: DescribeSlidesFunc | None = None,
    ) -> None:
        self._describe_slides_func = describe_slides_func

    async def process(
        self,
        source: SourceMaterial,
        *,
        router: ModelRouter | None = None,
    ) -> SourceDocument:
        if source.source_type != SourceType.PRESENTATION:
            raise UnsupportedFormatError(
                f"PresentationProcessor expects 'presentation', "
                f"got '{source.source_type}'"
            )

        path = Path(source.source_url)
        ext = path.suffix.lower()

        if ext not in SUPPORTED_EXTENSIONS:
            raise UnsupportedFormatError(
                f"Unsupported presentation format: {ext}. "
                f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
            )

        logger.info(
            "presentation_processing_start",
            source_url=source.source_url,
            format=ext,
        )

        if ext == ".pdf":
            chunks, page_count = await self._process_pdf(path)
        else:
            chunks, page_count = await self._process_pptx(path)

        logger.info(
            "presentation_processing_done",
            source_url=source.source_url,
            chunk_count=len(chunks),
        )

        return SourceDocument(
            source_type=SourceType.PRESENTATION,
            source_url=source.source_url,
            title=source.filename or path.stem,
            chunks=chunks,
            metadata={"page_count": page_count, "format": ext.lstrip(".")},
        )

    async def _process_pdf(
        self,
        path: Path,
    ) -> tuple[list[ContentChunk], int]:
        """Extract text (and optionally images) from PDF pages.

        Returns:
            Tuple of (chunks, page_count).
        """
        import fitz

        chunks: list[ContentChunk] = []
        doc = fitz.open(str(path))

        try:
            page_count = len(doc)

            for page_idx in range(page_count):
                page = doc[page_idx]
                slide_number = page_idx + 1
                text = page.get_text().strip()

                if text:
                    chunks.append(
                        ContentChunk(
                            chunk_type=ChunkType.SLIDE_TEXT,
                            text=text,
                            index=slide_number,
                            metadata={"slide_number": slide_number},
                        )
                    )
        finally:
            doc.close()

        # Optional: vision descriptions via injected heavy step
        if self._describe_slides_func is not None:
            descriptions = await self._describe_slides_func(
                str(path), DescribeSlidesParams()
            )
            for desc in descriptions:
                chunks.append(
                    ContentChunk(
                        chunk_type=ChunkType.SLIDE_DESCRIPTION,
                        text=desc.description,
                        index=desc.slide_number,
                        metadata={"slide_number": desc.slide_number},
                    )
                )

        return chunks, page_count

    async def _process_pptx(
        self,
        path: Path,
    ) -> tuple[list[ContentChunk], int]:
        """Extract text from PPTX slides.

        Returns:
            Tuple of (chunks, slide_count).
        """
        from pptx import Presentation

        chunks: list[ContentChunk] = []
        prs = Presentation(str(path))
        slides = list(prs.slides)

        for slide_idx, slide in enumerate(slides):
            slide_number = slide_idx + 1
            texts: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

            if texts:
                chunks.append(
                    ContentChunk(
                        chunk_type=ChunkType.SLIDE_TEXT,
                        text="\n".join(texts),
                        index=slide_number,
                        metadata={"slide_number": slide_number},
                    )
                )

        return chunks, len(slides)
