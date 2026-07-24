"""Tests for authored document API endpoints (tree-attached documents)."""

from __future__ import annotations

import io
import uuid
from datetime import UTC, datetime
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from fastapi import HTTPException
from httpx import ASGITransport, AsyncClient
from structlog.testing import capture_logs

from course_supporter.api.app import app
from course_supporter.api.deps import get_arq_redis, get_current_tenant, get_s3_client
from course_supporter.api.routes.documents import (
    _enforce_presentation_slide_cap,
    _intake_video_duration_sec,
    _presentation_slide_count,
)
from course_supporter.auth.context import TenantContext
from course_supporter.ingestion.base import ProcessingError, UnsupportedFormatError
from course_supporter.jobs.cancellation_service import JobCancellationService
from course_supporter.models.source import SourceType
from course_supporter.security.exceptions import ErrorCategory, SecurityRejectedError
from course_supporter.storage.authored_document_repository import (
    AuthoredDocumentRepository,
)
from course_supporter.storage.course_node_repository import CourseNodeRepository
from course_supporter.storage.database import get_session

STUB_TENANT = TenantContext(
    tenant_id=uuid.uuid4(),
    tenant_name="test-tenant",
    scopes=["prep", "check"],
    plan_id="basic",
    key_prefix="cs_test",
)

ENQUEUE_FUNC = "course_supporter.api.routes.documents.enqueue_ingestion"
PREP_ENQUEUE_FUNC = "course_supporter.api.routes.documents.enqueue_document_preparation"
RUN_STAGE1 = "course_supporter.api.routes.documents.run_stage1"
GET_MAX_SIZE = "course_supporter.api.routes.documents.get_max_size_for_extension"
PROBE_URL_FUNC = (
    "course_supporter.ingestion.video_pipeline.media.probe_intake_duration_sec"
)
PROBE_BYTES_FUNC = (
    "course_supporter.ingestion.video_pipeline.media.probe_intake_duration_from_bytes"
)


def _make_pdf_bytes(n_pages: int) -> bytes:
    """Synthesize a tiny valid PDF with ``n_pages`` blank pages (in-test)."""
    import fitz

    doc = fitz.open()
    for _ in range(n_pages):
        doc.new_page()
    data: bytes = doc.tobytes()
    doc.close()
    return data


def _make_pptx_bytes(n_slides: int) -> bytes:
    """Synthesize a tiny valid PPTX with ``n_slides`` blank slides (in-test)."""
    import io

    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(n_slides):
        prs.slides.add_slide(blank)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _mock_node(
    *,
    node_id: uuid.UUID | None = None,
    tenant_id: uuid.UUID | None = None,
) -> MagicMock:
    """Create a mock CourseNode with tenant_id."""
    node = MagicMock()
    node.id = node_id or uuid.uuid4()
    node.tenant_id = tenant_id or STUB_TENANT.tenant_id
    return node


def _mock_entry(
    *,
    entry_id: uuid.UUID | None = None,
    node_id: uuid.UUID | None = None,
    course_root_id: uuid.UUID | None = None,
    source_type: str = "text",
    material_role: str = "educational",
    task_type: str | None = None,
    source_url: str = "https://example.com/doc.md",
    filename: str | None = None,
    order: int = 0,
    state: str = "raw",
    error_message: str | None = None,
    job_id: uuid.UUID | None = None,
    slide_keys: list[str] | None = None,
) -> MagicMock:
    """Create a mock AuthoredDocument with ORM-compatible attributes."""
    entry = MagicMock()
    entry.id = entry_id or uuid.uuid4()
    entry.course_node_id = node_id or uuid.uuid4()
    # A-BE-2 (№21): explicit real UUID so MagicMock auto-attr does not feed the
    # response a MagicMock — the field is a required uuid.UUID (the ORM column is
    # NOT NULL), so a mock must always carry a concrete value.
    entry.course_root_id = course_root_id or uuid.uuid4()
    entry.source_type = source_type
    entry.material_role = material_role
    entry.task_type = task_type
    entry.source_url = source_url
    entry.filename = filename
    entry.language = None
    entry.order = order
    entry.state = state
    # L3: the response now carries a sibling ``processing_phase``. These tests
    # do not assert it — mirror the real derivation (error dominates; a
    # non-terminal mock defaults to the honest ``queued``).
    entry.processing_phase = (
        "error" if error_message else "ready" if state == "ready" else "queued"
    )
    entry.error_message = error_message
    entry.error_category = None
    entry.job_id = job_id
    # №21: explicit so MagicMock auto-attr does not feed
    # AuthoredDocumentResponse a MagicMock where a dict|None is expected.
    entry.file_roles = None
    entry.processing_estimate = None
    entry.deleted_at = None
    # Phase 6 T3: explicit so MagicMock auto-attr does not make the
    # delete-cleanup slide-key gathering branch spuriously truthy.
    entry.slide_keys = slide_keys
    entry.created_at = datetime.now(UTC)
    entry.updated_at = datetime.now(UTC)
    return entry


def _mock_job(job_id: uuid.UUID | None = None) -> MagicMock:
    """Create a mock Job returned by enqueue_ingestion."""
    job = MagicMock()
    job.id = job_id or uuid.uuid4()
    job.error_category = None
    return job


@pytest.fixture()
def mock_session() -> AsyncMock:
    session = AsyncMock()
    session.add = MagicMock()
    return session


@pytest.fixture()
def mock_arq() -> MagicMock:
    return MagicMock()


@pytest.fixture()
def mock_s3() -> AsyncMock:
    s3 = AsyncMock()
    s3.upload_smart = AsyncMock(
        return_value=("http://localhost:9000/course-materials/key/file.pdf", 1024)
    )
    s3.extract_key = MagicMock(return_value=None)
    return s3


@pytest.fixture()
def node_id() -> uuid.UUID:
    return uuid.uuid4()


@pytest.fixture()
async def client(
    mock_session: AsyncMock, mock_arq: MagicMock, mock_s3: AsyncMock
) -> AsyncClient:
    app.dependency_overrides[get_session] = lambda: mock_session
    app.dependency_overrides[get_current_tenant] = lambda: STUB_TENANT
    app.dependency_overrides[get_arq_redis] = lambda: mock_arq
    app.dependency_overrides[get_s3_client] = lambda: mock_s3
    async with AsyncClient(
        transport=ASGITransport(app=app),
        base_url="http://test",
    ) as ac:
        yield ac  # type: ignore[misc]
    app.dependency_overrides.clear()


@pytest.fixture(autouse=True)
def _reset_rate_limiter() -> None:
    """Isolate the process-global rate limiter across tests.

    Every api test authenticates as the same stub tenant → one rate-limit
    bucket; without a per-test reset a long test file eventually trips its own
    429 (the sliding window is 60s, far longer than the whole run). Not a
    behaviour under test here — rate limiting has its own suite.
    """
    from course_supporter.auth.scopes import rate_limiter

    rate_limiter.reset()


class TestCreateDocument:
    """POST /api/v1/nodes/{nid}/documents"""

    async def test_returns_201_with_url(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Successful document creation with URL returns 201 with job_id."""
        entry = _mock_entry(node_id=node_id)
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "text",
                    "source_url": "https://example.com/doc.md",
                },
            )
        assert resp.status_code == 201
        data = resp.json()
        assert data["id"] == str(entry.id)
        assert data["job_id"] == str(job.id)
        assert data["source_type"] == "text"

    async def test_returns_201_with_file_upload(
        self,
        client: AsyncClient,
        node_id: uuid.UUID,
        mock_s3: AsyncMock,
    ) -> None:
        """Successful file upload returns 201.

        Phase 1.2 C3 wired Stage 1 magic validation at the multipart
        callsite. Fixture bytes ``b"PDF content"`` are not a valid
        PDF magic header — patch ``run_stage1`` to bypass the magic
        check so the test continues to exercise the happy-path
        wiring (Stage 1 ok → S3 upload → repo create → 201).
        Phase 0.6 sealed-library tests cover magic validation at the
        unit tier; here we test the route-level S3 + repo flow.
        """
        entry = _mock_entry(
            node_id=node_id,
            source_type="presentation",
            source_url="http://localhost:9000/course-materials/key/file.pdf",
            filename="slides.pdf",
        )
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
            patch(
                "course_supporter.api.routes.documents.run_stage1",
                return_value=None,
            ),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={"source_type": "presentation"},
                files={
                    "file": (
                        "slides.pdf",
                        io.BytesIO(b"PDF content"),
                        "application/pdf",
                    )
                },
            )
        assert resp.status_code == 201
        mock_s3.upload_smart.assert_awaited_once()

    async def test_invalid_source_type_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Invalid source_type is rejected by validation."""
        resp = await client.post(
            f"/api/v1/nodes/{node_id}/documents",
            data={
                "source_type": "invalid",
                "source_url": "https://example.com/doc.md",
            },
        )
        assert resp.status_code == 422

    async def test_no_url_no_file_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Neither URL nor file provided returns 422."""
        resp = await client.post(
            f"/api/v1/nodes/{node_id}/documents",
            data={"source_type": "text"},
        )
        assert resp.status_code == 422
        assert "Either source_url or file" in resp.json()["detail"]

    async def test_web_rejects_file_upload(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """source_type 'web' does not accept file uploads."""
        resp = await client.post(
            f"/api/v1/nodes/{node_id}/documents",
            data={"source_type": "web"},
            files={
                "file": ("page.html", io.BytesIO(b"<html>"), "text/html"),
            },
        )
        assert resp.status_code == 422
        assert "does not accept file uploads" in resp.json()["detail"]

    async def test_forbidden_extension_returns_400(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Extension outside ``AUTHORED_POLICY`` whitelist returns 400.

        Phase 1.2 C3 — KD14 envelope replaces legacy 422 + freeform
        ``detail`` string. Source-type-segmented rejection (e.g.
        ``.pdf`` for ``video``) was removed; the policy is flat. Use
        ``.exe`` (always-rejected) to exercise the FORBIDDEN_TYPE
        wiring at the multipart create_document endpoint.
        """
        resp = await client.post(
            f"/api/v1/nodes/{node_id}/documents",
            data={"source_type": "video"},
            files={
                "file": (
                    "payload.exe",
                    io.BytesIO(b"binary"),
                    "application/octet-stream",
                ),
            },
        )
        assert resp.status_code == 400
        detail = resp.json()["detail"]
        assert detail["code"] == "SECURITY_REJECTED"
        assert detail["category"] == "forbidden_type"
        assert "details" in detail

    async def test_node_not_found_returns_404(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Non-existent node returns 404."""
        with patch.object(CourseNodeRepository, "get_by_id", return_value=None):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "text",
                    "source_url": "https://example.com/doc.md",
                },
            )
        assert resp.status_code == 404
        assert resp.json()["detail"] == "Node not found"

    async def test_node_wrong_tenant_returns_404(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Node belonging to another tenant returns 404."""
        other_tenant = uuid.uuid4()
        with patch.object(
            CourseNodeRepository,
            "get_by_id",
            return_value=_mock_node(node_id=node_id, tenant_id=other_tenant),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "text",
                    "source_url": "https://example.com/doc.md",
                },
            )
        assert resp.status_code == 404

    async def test_with_filename_override(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Creation with filename override includes it in response."""
        entry = _mock_entry(node_id=node_id, filename="notes.md")
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "text",
                    "source_url": "https://example.com/notes.md",
                    "filename": "notes.md",
                },
            )
        assert resp.status_code == 201
        assert resp.json()["filename"] == "notes.md"

    async def test_language_639_1_normalized_to_639_3(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Form-input 639-1 ``uk`` → canonical 639-3 ``ukr`` (Task 2.4.13)."""
        captured: dict[str, object] = {}

        async def fake_create(self: object, **kwargs: object) -> MagicMock:
            captured.update(kwargs)
            entry = _mock_entry(node_id=node_id)
            entry.language = kwargs.get("language")
            return entry

        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", new=fake_create),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "text",
                    "source_url": "https://example.com/doc.md",
                    "language": "uk",
                },
            )
        assert resp.status_code == 201
        assert captured["language"] == "ukr"
        assert resp.json()["language"] == "ukr"

    async def test_language_invalid_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Form-input language not in the whitelist returns 422, not 500."""
        with patch.object(
            CourseNodeRepository,
            "get_by_id",
            return_value=_mock_node(node_id=node_id),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "text",
                    "source_url": "https://example.com/doc.md",
                    "language": "xyz",
                },
            )
        assert resp.status_code == 422

    async def test_language_absent_succeeds(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Optional ``language`` field — absence is accepted."""
        entry = _mock_entry(node_id=node_id)
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "text",
                    "source_url": "https://example.com/doc.md",
                },
            )
        assert resp.status_code == 201


class TestCodeUploadRouting:
    """task-code-materials commit 3 — multipart light-gate for source_type=code.

    Code uploads (single source file or a .zip project archive) bypass the
    fail-closed ``run_stage1`` (base-attach precedent): extension routing +
    size cap only. Content safety is the async Stage-2 check; the archive
    sandbox is ``extract_archive_safely`` inside the CodeProcessor worker.
    """

    async def test_code_single_file_201_skips_stage1(
        self,
        client: AsyncClient,
        node_id: uuid.UUID,
        mock_s3: AsyncMock,
    ) -> None:
        entry = _mock_entry(
            node_id=node_id,
            source_type="code",
            source_url="http://localhost:9000/course-materials/key/script.py",
            filename="script.py",
        )
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(
                PREP_ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job
            ) as prep_mock,
            patch(ENQUEUE_FUNC, new_callable=AsyncMock) as ingest_mock,
            patch(RUN_STAGE1) as stage1_mock,
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={"source_type": "code"},
                files={
                    "file": (
                        "script.py",
                        io.BytesIO(b'print("hello")\n'),
                        "text/x-python",
                    )
                },
            )
        assert resp.status_code == 201
        stage1_mock.assert_not_called()
        # BE8: a code upload routes to DOCUMENT_PREPARATION, never straight to
        # ingestion (the author confirms file roles before the expensive pass).
        prep_mock.assert_awaited_once()
        ingest_mock.assert_not_awaited()
        mock_s3.upload_smart.assert_awaited_once()

    async def test_code_zip_archive_201_skips_stage1(
        self,
        client: AsyncClient,
        node_id: uuid.UUID,
        mock_s3: AsyncMock,
    ) -> None:
        entry = _mock_entry(
            node_id=node_id,
            source_type="code",
            source_url="http://localhost:9000/course-materials/key/project.zip",
            filename="project.zip",
        )
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(
                PREP_ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job
            ) as prep_mock,
            patch(ENQUEUE_FUNC, new_callable=AsyncMock) as ingest_mock,
            patch(RUN_STAGE1) as stage1_mock,
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={"source_type": "code"},
                files={
                    "file": (
                        "project.zip",
                        io.BytesIO(b"PK\x03\x04 stub"),
                        "application/zip",
                    )
                },
            )
        assert resp.status_code == 201
        stage1_mock.assert_not_called()
        # BE8: a code (zip) upload also routes to DOCUMENT_PREPARATION.
        prep_mock.assert_awaited_once()
        ingest_mock.assert_not_awaited()

    async def test_code_non_code_extension_400_forbidden(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        # The light gate still rejects extensions outside CODE_EXTENSIONS/zip.
        resp = await client.post(
            f"/api/v1/nodes/{node_id}/documents",
            data={"source_type": "code"},
            files={
                "file": ("movie.mp4", io.BytesIO(b"\x00\x00"), "video/mp4"),
            },
        )
        assert resp.status_code == 400
        detail = resp.json()["detail"]
        assert detail["code"] == "SECURITY_REJECTED"
        assert detail["category"] == "forbidden_type"

    async def test_zip_with_non_code_source_type_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        # zip exists in the authored whitelist ONLY for code — a zip declared
        # as any other source_type has no processor; deterministic fast-fail.
        resp = await client.post(
            f"/api/v1/nodes/{node_id}/documents",
            data={"source_type": "text"},
            files={
                "file": ("bundle.zip", io.BytesIO(b"PK\x03\x04"), "application/zip"),
            },
        )
        assert resp.status_code == 422
        assert resp.json()["detail"]["code"] == "ARCHIVE_REQUIRES_CODE"

    async def test_code_requires_file_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        resp = await client.post(
            f"/api/v1/nodes/{node_id}/documents",
            data={
                "source_type": "code",
                "source_url": "https://example.com/script.py",
            },
        )
        assert resp.status_code == 422
        assert "requires a file upload" in resp.json()["detail"]


class TestPresentationSlideCountInspection:
    """KD-2.3-M -- _presentation_slide_count helper (HTTP-side, <=200 ms)."""

    def test_pdf_counts_pages(self) -> None:
        assert _presentation_slide_count("pdf", _make_pdf_bytes(7)) == 7

    def test_pptx_counts_slides(self) -> None:
        assert _presentation_slide_count("pptx", _make_pptx_bytes(4)) == 4

    def test_ppt_skips_returns_none(self) -> None:
        # Legacy binary -- python-pptx cannot read it; defer to worker (CA-3).
        assert _presentation_slide_count("ppt", b"\xd0\xcf\x11\xe0legacy") is None

    def test_non_presentation_extension_returns_none(self) -> None:
        assert _presentation_slide_count("md", b"# hello") is None

    def test_skip_on_corrupt_pdf(self) -> None:
        # Magic-valid (%PDF) but structurally broken -> fitz.FileDataError ->
        # skip with a warning (defer to worker); no exception escapes (CA-4).
        with capture_logs() as logs:
            result = _presentation_slide_count("pdf", b"%PDF-1.4 dummy")
        assert result is None
        assert any(e["event"] == "slide_count_inspect_failed" for e in logs)

    def test_skip_on_corrupt_pptx(self) -> None:
        # Non-ZIP garbage -> zipfile.BadZipFile -> skip with a warning (CA-4).
        with capture_logs() as logs:
            result = _presentation_slide_count("pptx", b"not a pptx at all")
        assert result is None
        assert any(e["event"] == "slide_count_inspect_failed" for e in logs)


class TestEnforcePresentationSlideCap:
    """Guard semantics for _enforce_presentation_slide_cap (KD-2.3-M + B guard)."""

    def test_raises_over_limit_for_presentation(self) -> None:
        with pytest.raises(SecurityRejectedError) as exc_info:
            _enforce_presentation_slide_cap(
                SourceType.PRESENTATION, "pdf", _make_pdf_bytes(101)
            )
        assert exc_info.value.category == ErrorCategory.SLIDE_COUNT_LIMIT

    def test_passes_at_limit(self) -> None:
        # Exactly 100 is allowed; only > 100 rejects.
        _enforce_presentation_slide_cap(
            SourceType.PRESENTATION, "pdf", _make_pdf_bytes(100)
        )

    def test_noop_for_non_presentation_source_type(self) -> None:
        # B guard: a long text-source PDF is NOT a presentation -> no check.
        _enforce_presentation_slide_cap(SourceType.TEXT, "pdf", _make_pdf_bytes(150))

    def test_noop_for_ppt(self) -> None:
        _enforce_presentation_slide_cap(
            SourceType.PRESENTATION, "ppt", b"\xd0\xcf\x11\xe0"
        )


class TestPresentationSlideCountRoute:
    """KD-2.3-M -- slide-count gate wired into the upload endpoints (gate 23).

    FORBIDDEN_TYPE prong: ``test_forbidden_type_for_non_whitelisted_extension``
    below. NOTE the sealed gate-23 example uses ``.docx``, but docx IS in
    AUTHORED_POLICY.allowed_extensions (accepted, 201) -- so a genuinely
    non-whitelisted extension exercises the FORBIDDEN_TYPE wiring instead.
    """

    async def test_pdf_over_limit_returns_slide_count_limit(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(RUN_STAGE1, return_value=None),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={"source_type": "presentation"},
                files={
                    "file": (
                        "deck.pdf",
                        io.BytesIO(_make_pdf_bytes(101)),
                        "application/pdf",
                    )
                },
            )
        assert resp.status_code == 400
        detail = resp.json()["detail"]
        assert detail["code"] == "SECURITY_REJECTED"
        assert detail["category"] == "slide_count_limit"

    async def test_pdf_under_limit_passes(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        entry = _mock_entry(node_id=node_id, source_type="presentation")
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
            patch(RUN_STAGE1, return_value=None),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={"source_type": "presentation"},
                files={
                    "file": (
                        "deck.pdf",
                        io.BytesIO(_make_pdf_bytes(5)),
                        "application/pdf",
                    )
                },
            )
        assert resp.status_code == 201

    async def test_ppt_skips_slide_count(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        # .ppt skips the HTTP gate (CA-3 -> worker enforces): 201 even though
        # the file is a presentation, because python-pptx cannot read .ppt.
        entry = _mock_entry(node_id=node_id, source_type="presentation")
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
            patch(RUN_STAGE1, return_value=None),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={"source_type": "presentation"},
                files={
                    "file": (
                        "legacy.ppt",
                        io.BytesIO(b"\xd0\xcf\x11\xe0legacy-binary"),
                        "application/vnd.ms-powerpoint",
                    )
                },
            )
        assert resp.status_code == 201

    async def test_text_source_pdf_skips_gate(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        # B guard: a 150-page PDF uploaded as source_type=text is NOT subject
        # to the presentation slide-count gate.
        entry = _mock_entry(node_id=node_id, source_type="text")
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
            patch(RUN_STAGE1, return_value=None),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={"source_type": "text"},
                files={
                    "file": (
                        "notes.pdf",
                        io.BytesIO(_make_pdf_bytes(150)),
                        "application/pdf",
                    )
                },
            )
        assert resp.status_code == 201

    async def test_oversize_presentation_returns_size_limit(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        # SIZE_LIMIT prong of gate 23. The 50 MB literal is unit-tested in
        # test_policies; here the resolver is patched to a tiny cap to
        # exercise the route wiring without a 50 MB fixture (ratified strategy).
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(GET_MAX_SIZE, return_value=5),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={"source_type": "presentation"},
                files={
                    "file": (
                        "deck.pdf",
                        io.BytesIO(_make_pdf_bytes(2)),
                        "application/pdf",
                    )
                },
            )
        assert resp.status_code == 400
        assert resp.json()["detail"]["category"] == "size_limit"

    async def test_forbidden_type_for_non_whitelisted_extension(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        # Gate 23 FORBIDDEN_TYPE prong via a genuinely non-whitelisted ext
        # (real run_stage1 rejects on the extension whitelist).
        resp = await client.post(
            f"/api/v1/nodes/{node_id}/documents",
            data={"source_type": "presentation"},
            files={
                "file": (
                    "payload.exe",
                    io.BytesIO(b"binary"),
                    "application/octet-stream",
                )
            },
        )
        assert resp.status_code == 400
        assert resp.json()["detail"]["category"] == "forbidden_type"

    async def test_confirm_upload_pptx_over_limit(
        self, client: AsyncClient, node_id: uuid.UUID, mock_s3: AsyncMock
    ) -> None:
        # D: the presigned confirm path enforces the slide-count gate too.
        pptx_bytes = _make_pptx_bytes(110)
        key = f"tenants/{STUB_TENANT.tenant_id}/nodes/{node_id}/abc/deck.pptx"
        mock_s3.head_object = AsyncMock(return_value={"ContentLength": len(pptx_bytes)})
        mock_s3.get_object = AsyncMock(return_value=pptx_bytes)
        mock_s3.get_object_url = MagicMock(return_value="http://localhost:9000/x")
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(RUN_STAGE1, return_value=None),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents/confirm-upload",
                json={"key": key, "source_type": "presentation"},
            )
        assert resp.status_code == 400
        assert resp.json()["detail"]["category"] == "slide_count_limit"


class TestCreateDocumentOpenAPISpec:
    """Regression: ensure create_document endpoint accepts multipart/form-data."""

    async def test_openapi_content_type_is_multipart(self, client: AsyncClient) -> None:
        """File upload endpoint must declare multipart/form-data content type."""
        resp = await client.get("/openapi.json")
        schema = resp.json()
        path = "/api/v1/nodes/{node_id}/documents"
        content_types = list(
            schema["paths"][path]["post"]["requestBody"]["content"].keys()
        )
        assert "multipart/form-data" in content_types


class TestListDocuments:
    """GET /api/v1/nodes/{nid}/documents"""

    async def test_returns_list(self, client: AsyncClient, node_id: uuid.UUID) -> None:
        """Returns list of documents for the node."""
        entries = [
            _mock_entry(node_id=node_id, order=0),
            _mock_entry(node_id=node_id, order=1, source_type="video"),
        ]
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(
                AuthoredDocumentRepository, "get_for_node", return_value=entries
            ),
        ):
            resp = await client.get(f"/api/v1/nodes/{node_id}/documents")
        assert resp.status_code == 200
        data = resp.json()
        assert len(data) == 2
        assert data[0]["order"] == 0
        assert data[1]["source_type"] == "video"

    async def test_empty_list(self, client: AsyncClient, node_id: uuid.UUID) -> None:
        """Returns empty list when node has no documents."""
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "get_for_node", return_value=[]),
        ):
            resp = await client.get(f"/api/v1/nodes/{node_id}/documents")
        assert resp.status_code == 200
        assert resp.json() == []

    async def test_node_not_found_returns_404(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Non-existent node returns 404."""
        with patch.object(CourseNodeRepository, "get_by_id", return_value=None):
            resp = await client.get(f"/api/v1/nodes/{node_id}/documents")
        assert resp.status_code == 404


class TestGetDocument:
    """GET /api/v1/documents/{did}"""

    async def test_returns_entry(self, client: AsyncClient, node_id: uuid.UUID) -> None:
        """Returns single authored document."""
        entry = _mock_entry(node_id=node_id, state="ready")
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.get(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 200
        data = resp.json()
        assert data["id"] == str(entry.id)
        assert data["state"] == "ready"

    async def test_returns_course_root_id(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """A-BE-2 (№21): the response projects the stored course_root_id."""
        root_id = uuid.uuid4()
        entry = _mock_entry(node_id=node_id, course_root_id=root_id, state="ready")
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.get(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 200
        assert resp.json()["course_root_id"] == str(root_id)

    async def test_not_found_returns_404(self, client: AsyncClient) -> None:
        """Non-existent document returns 404."""
        with patch.object(AuthoredDocumentRepository, "get_by_id", return_value=None):
            resp = await client.get(f"/api/v1/documents/{uuid.uuid4()}")
        assert resp.status_code == 404

    async def test_wrong_tenant_returns_404(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Document belonging to another tenant returns 404."""
        entry = _mock_entry(node_id=node_id)
        other_tenant = uuid.uuid4()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id, tenant_id=other_tenant),
            ),
        ):
            resp = await client.get(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 404


class TestDeleteDocument:
    """DELETE /api/v1/documents/{did} — Phase 1 commit (l) KD3 cascade
    soft-delete + s3_cleanup orchestration.

    Mirrors commit (k)'s ``TestDeleteNode`` shape with the cascade
    rooted at AuthoredDocument instead of CourseNode. The handler now:
    (1) extracts the S3 key from ``document.source_url`` BEFORE
    cascade fires (cascade scrub clears ``source_url`` to ``""``),
    (2) issues ``CascadeDeleteService.soft_delete_with_cascade`` which
    auto-dispatches ``__scrub_callable__`` per victim type
    (AuthoredDocument scrubbed; DocumentSummary + DocumentSegment
    descendants no-op in Phase 1 per Amendment 16), and (3) hands off
    to ``enqueue_s3_cleanup`` which owns the QQ5 commit boundary.
    """

    async def test_returns_204_when_no_s3_key(
        self, client: AsyncClient, node_id: uuid.UUID, mock_s3: AsyncMock
    ) -> None:
        """External URL (extract_key returns None) — handler short-circuits
        to direct ``session.commit()``; no enqueue, still 204.
        """
        entry = _mock_entry(node_id=node_id, source_url="https://example.com/doc.md")
        mock_s3.extract_key = MagicMock(return_value=None)
        cascade_mock = AsyncMock()
        enqueue_mock = AsyncMock()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(
                "course_supporter.storage.cascade.CascadeDeleteService."
                "soft_delete_with_cascade",
                cascade_mock,
            ),
            patch(
                "course_supporter.api.routes.documents.enqueue_s3_cleanup",
                enqueue_mock,
            ),
        ):
            resp = await client.delete(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 204
        cascade_mock.assert_awaited_once()
        enqueue_mock.assert_not_awaited()

    async def test_not_found_returns_404(self, client: AsyncClient) -> None:
        """Non-existent document returns 404."""
        with patch.object(AuthoredDocumentRepository, "get_by_id", return_value=None):
            resp = await client.delete(f"/api/v1/documents/{uuid.uuid4()}")
        assert resp.status_code == 404

    async def test_wrong_tenant_returns_404(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Document belonging to another tenant returns 404."""
        entry = _mock_entry(node_id=node_id)
        other_tenant = uuid.uuid4()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id, tenant_id=other_tenant),
            ),
        ):
            resp = await client.delete(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 404

    async def test_collects_key_before_cascade_and_enqueues_cleanup(
        self,
        client: AsyncClient,
        node_id: uuid.UUID,
        mock_s3: AsyncMock,
    ) -> None:
        """File key extracted from ``document.source_url`` and forwarded
        to ``enqueue_s3_cleanup`` along with tenant/course_node anchors.
        Locks the QQ5 ordering (collect → cascade → enqueue) at the
        unit level — the cascade engine integration test in
        ``tests/storage/test_cascade_invalidation.py`` locks the
        scrub-then-collect-impossible failure mode.
        """
        entry = _mock_entry(
            node_id=node_id,
            source_url="http://localhost:9000/bucket/tenants/t/file.pdf",
        )
        mock_s3.extract_key = MagicMock(return_value="tenants/t/file.pdf")
        cascade_mock = AsyncMock()
        enqueue_mock = AsyncMock()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(
                "course_supporter.storage.cascade.CascadeDeleteService."
                "soft_delete_with_cascade",
                cascade_mock,
            ),
            patch(
                "course_supporter.api.routes.documents.enqueue_s3_cleanup",
                enqueue_mock,
            ),
        ):
            resp = await client.delete(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 204
        cascade_mock.assert_awaited_once()
        enqueue_mock.assert_awaited_once()
        kwargs = enqueue_mock.call_args.kwargs
        assert kwargs["file_keys"] == ["tenants/t/file.pdf"]
        assert kwargs["course_node_id"] == node_id
        assert kwargs["tenant_id"] == STUB_TENANT.tenant_id

    async def test_gathers_slide_keys_into_cleanup(
        self,
        client: AsyncClient,
        node_id: uuid.UUID,
        mock_s3: AsyncMock,
    ) -> None:
        """Persisted per-slide WebP keys (Phase 6 T3) are scrubbed too.

        A presentation's ``slide_keys`` are S3 objects distinct from
        ``source_url``; the handler gathers them (after the source key,
        before cascade) so ``enqueue_s3_cleanup`` removes every WebP slide
        alongside the original deck.
        """
        entry = _mock_entry(
            node_id=node_id,
            source_type="presentation",
            source_url="http://localhost:9000/bucket/tenants/t/deck.pptx",
            slide_keys=[
                "tenants/t/nodes/n/slides/d/0001.webp",
                "tenants/t/nodes/n/slides/d/0002.webp",
            ],
        )
        mock_s3.extract_key = MagicMock(return_value="tenants/t/deck.pptx")
        cascade_mock = AsyncMock()
        enqueue_mock = AsyncMock()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(
                "course_supporter.storage.cascade.CascadeDeleteService."
                "soft_delete_with_cascade",
                cascade_mock,
            ),
            patch(
                "course_supporter.api.routes.documents.enqueue_s3_cleanup",
                enqueue_mock,
            ),
        ):
            resp = await client.delete(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 204
        enqueue_mock.assert_awaited_once()
        # Source deck key first, then the two slide keys — ordered, all scrubbed.
        assert enqueue_mock.call_args.kwargs["file_keys"] == [
            "tenants/t/deck.pptx",
            "tenants/t/nodes/n/slides/d/0001.webp",
            "tenants/t/nodes/n/slides/d/0002.webp",
        ]

    async def test_no_enqueue_when_extract_key_returns_none(
        self,
        client: AsyncClient,
        node_id: uuid.UUID,
        mock_s3: AsyncMock,
    ) -> None:
        """External URL yields ``None`` from ``extract_key`` — handler
        skips ``enqueue_s3_cleanup`` and commits the cascade directly.
        Avoids creating a wasteful Job row + ARQ task with empty payload.
        """
        entry = _mock_entry(node_id=node_id, source_url="https://example.com/video.mp4")
        mock_s3.extract_key = MagicMock(return_value=None)
        cascade_mock = AsyncMock()
        enqueue_mock = AsyncMock()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(
                "course_supporter.storage.cascade.CascadeDeleteService."
                "soft_delete_with_cascade",
                cascade_mock,
            ),
            patch(
                "course_supporter.api.routes.documents.enqueue_s3_cleanup",
                enqueue_mock,
            ),
        ):
            resp = await client.delete(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 204
        cascade_mock.assert_awaited_once()
        enqueue_mock.assert_not_awaited()

    async def test_passes_cancel_hook_direct_bind_no_augmentation(
        self,
        client: AsyncClient,
        node_id: uuid.UUID,
        mock_s3: AsyncMock,
    ) -> None:
        """L1b — ``on_cancel_jobs`` is the DIRECT-bound
        :meth:`JobCancellationService.cancel_jobs_for_entities`, not a closure.
        JCS now keys on ``Job.subject_id IN victim_ids`` and the document's own
        id is in the cascade victim set, so the document's job cancels precisely
        — no ``course_node_id`` augmentation (which caused the F3 over-cancel).
        """
        entry = _mock_entry(
            node_id=node_id,
            source_url="https://example.com/doc.md",
        )
        mock_s3.extract_key = MagicMock(return_value=None)
        cascade_mock = AsyncMock()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(
                "course_supporter.storage.cascade.CascadeDeleteService."
                "soft_delete_with_cascade",
                cascade_mock,
            ),
        ):
            resp = await client.delete(f"/api/v1/documents/{entry.id}")
            assert resp.status_code == 204
            cascade_mock.assert_awaited_once()
            on_cancel_jobs = cascade_mock.call_args.kwargs.get("on_cancel_jobs")
            assert on_cancel_jobs is not None, "on_cancel_jobs missing — KD13 gap"
            # Direct bind (the nodes.py pattern) — a bound method whose __func__
            # is the JCS method itself, NOT a wrapping closure.
            assert (
                getattr(on_cancel_jobs, "__func__", None)
                is JobCancellationService.cancel_jobs_for_entities
            ), (
                "delete_document must direct-bind cancel_jobs_for_entities "
                "(L1b — no course_node_id augmentation)"
            )

    async def test_passes_on_invalidate_hashes_hook(
        self,
        client: AsyncClient,
        node_id: uuid.UUID,
        mock_s3: AsyncMock,
    ) -> None:
        """Hotfix-5 contract — closes tangential gap from commit (l)
        where ``on_invalidate_hashes`` was omitted entirely from
        ``delete_document`` (mirrors of this hook are present in
        ``delete_node`` (commit (k)) and ``delete_file`` (commit (m))).
        Regression guard: a future revert to a 2-arg cascade call would
        re-introduce stale parent-CourseNode content_hash bug.
        """
        entry = _mock_entry(
            node_id=node_id,
            source_url="https://example.com/doc.md",
        )
        mock_s3.extract_key = MagicMock(return_value=None)
        cascade_mock = AsyncMock()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(
                "course_supporter.storage.cascade.CascadeDeleteService."
                "soft_delete_with_cascade",
                cascade_mock,
            ),
        ):
            resp = await client.delete(f"/api/v1/documents/{entry.id}")
        assert resp.status_code == 204
        cascade_mock.assert_awaited_once()
        on_invalidate_hashes = cascade_mock.call_args.kwargs.get("on_invalidate_hashes")
        assert on_invalidate_hashes is not None, (
            "on_invalidate_hashes missing — tangential gap from commit "
            "(l) regressed; parent-CourseNode content_hash recompute "
            "would treat the document as still present"
        )


class TestRetryDocument:
    """POST /api/v1/documents/{did}/retry"""

    async def test_returns_200_with_new_job(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Successful retry returns 200 with new job_id."""
        entry = _mock_entry(
            node_id=node_id,
            state="error",
            error_message="Processing failed",
        )
        job = _mock_job()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            # L1b: retry supersedes any in-flight job of the subject before
            # enqueue — mocked here (the SQL path is exercised in the DB tier).
            patch.object(
                JobCancellationService,
                "cancel_jobs_for_entities",
                new_callable=AsyncMock,
            ),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
        ):
            resp = await client.post(f"/api/v1/documents/{entry.id}/retry")
        assert resp.status_code == 200
        data = resp.json()
        assert data["job_id"] == str(job.id)

    async def test_code_retry_routes_to_preparation(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """BE8: retry of a code doc re-runs DOCUMENT_PREPARATION (decision 19),
        never straight-to-processing; non-code retries stay on ingestion (I6)."""
        entry = _mock_entry(
            node_id=node_id,
            source_type="code",
            state="error",
            error_message="Processing failed",
        )
        job = _mock_job()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(
                JobCancellationService,
                "cancel_jobs_for_entities",
                new_callable=AsyncMock,
            ),
            patch(
                PREP_ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job
            ) as prep_mock,
            patch(ENQUEUE_FUNC, new_callable=AsyncMock) as ingest_mock,
        ):
            resp = await client.post(f"/api/v1/documents/{entry.id}/retry")
        assert resp.status_code == 200
        prep_mock.assert_awaited_once()
        ingest_mock.assert_not_awaited()

    async def test_non_error_state_returns_409(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Retry on non-error document returns 409."""
        entry = _mock_entry(node_id=node_id, state="ready")
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.post(f"/api/v1/documents/{entry.id}/retry")
        assert resp.status_code == 409
        assert "ready" in resp.json()["detail"]

    async def test_pending_state_returns_409(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Retry on pending document returns 409."""
        entry = _mock_entry(node_id=node_id, state="pending")
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.post(f"/api/v1/documents/{entry.id}/retry")
        assert resp.status_code == 409

    async def test_not_found_returns_404(self, client: AsyncClient) -> None:
        """Non-existent document returns 404."""
        with patch.object(AuthoredDocumentRepository, "get_by_id", return_value=None):
            resp = await client.post(f"/api/v1/documents/{uuid.uuid4()}/retry")
        assert resp.status_code == 404

    async def test_returns_410_when_soft_deleted(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """QQ6: retry on a soft-deleted document yields HTTP 410 Gone.

        Fires BEFORE the state-check branch — even an ``error``-state
        soft-deleted row returns 410, never 409. The row's been
        logically removed; re-ingestion is not a recovery path.
        """
        entry = _mock_entry(
            node_id=node_id,
            state="error",
            error_message="Processing failed",
        )
        entry.deleted_at = datetime.now(UTC)
        enqueue_mock = AsyncMock()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(ENQUEUE_FUNC, enqueue_mock),
        ):
            resp = await client.post(f"/api/v1/documents/{entry.id}/retry")
        assert resp.status_code == 410
        assert "deleted" in resp.json()["detail"].lower()
        # Retry orchestration must NOT fire when soft-deleted.
        enqueue_mock.assert_not_awaited()


class TestUpdateDocument:
    """PATCH /api/v1/documents/{did}"""

    async def test_updates_role_to_methodological(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Successful role update returns 200 with new role."""
        entry = _mock_entry(node_id=node_id, material_role="educational")
        updated = _mock_entry(node_id=node_id, material_role="methodological")
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(
                AuthoredDocumentRepository,
                "update_material_role",
                return_value=updated,
            ),
        ):
            resp = await client.patch(
                f"/api/v1/documents/{entry.id}",
                json={"material_role": "methodological"},
            )
        assert resp.status_code == 200
        assert resp.json()["material_role"] == "methodological"

    async def test_updates_role_to_educational(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Toggle back to educational."""
        entry = _mock_entry(node_id=node_id, material_role="methodological")
        updated = _mock_entry(node_id=node_id, material_role="educational")
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(
                AuthoredDocumentRepository,
                "update_material_role",
                return_value=updated,
            ),
        ):
            resp = await client.patch(
                f"/api/v1/documents/{entry.id}",
                json={"material_role": "educational"},
            )
        assert resp.status_code == 200
        assert resp.json()["material_role"] == "educational"

    async def test_invalid_role_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Invalid material_role value returns 422."""
        resp = await client.patch(
            f"/api/v1/documents/{uuid.uuid4()}",
            json={"material_role": "invalid"},
        )
        assert resp.status_code == 422

    async def test_set_task_type(self, client: AsyncClient, node_id: uuid.UUID) -> None:
        """PATCH sets task_type without changing role."""
        entry = _mock_entry(node_id=node_id, material_role="educational")
        updated = _mock_entry(
            node_id=node_id,
            material_role="educational",
            task_type="short_task",
        )
        update_task_mock = AsyncMock(return_value=updated)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(
                AuthoredDocumentRepository,
                "update_task_type",
                update_task_mock,
            ),
        ):
            resp = await client.patch(
                f"/api/v1/documents/{entry.id}",
                json={"task_type": "short_task"},
            )
        assert resp.status_code == 200
        assert resp.json()["task_type"] == "short_task"
        # material_role update should NOT have been called since it wasn't in body
        update_task_mock.assert_awaited_once()

    async def test_clear_task_type(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """PATCH with task_type: null clears the task flag."""
        entry = _mock_entry(
            node_id=node_id,
            material_role="educational",
            task_type="task",
        )
        updated = _mock_entry(
            node_id=node_id,
            material_role="educational",
            task_type=None,
        )
        update_task_mock = AsyncMock(return_value=updated)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(
                AuthoredDocumentRepository,
                "update_task_type",
                update_task_mock,
            ),
        ):
            resp = await client.patch(
                f"/api/v1/documents/{entry.id}",
                json={"task_type": None},
            )
        assert resp.status_code == 200
        assert resp.json()["task_type"] is None
        update_task_mock.assert_awaited_once()

    async def test_empty_body_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """PATCH with no fields returns 422."""
        entry = _mock_entry(node_id=node_id)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.patch(
                f"/api/v1/documents/{entry.id}",
                json={},
            )
        assert resp.status_code == 422

    async def test_invalid_task_type_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """PATCH rejects task_type outside the taxonomy."""
        resp = await client.patch(
            f"/api/v1/documents/{uuid.uuid4()}",
            json={"task_type": "essay"},
        )
        assert resp.status_code == 422

    async def test_not_found_returns_404(self, client: AsyncClient) -> None:
        """Non-existent document returns 404."""
        with patch.object(AuthoredDocumentRepository, "get_by_id", return_value=None):
            resp = await client.patch(
                f"/api/v1/documents/{uuid.uuid4()}",
                json={"material_role": "methodological"},
            )
        assert resp.status_code == 404

    async def test_wrong_tenant_returns_404(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """Document belonging to another tenant returns 404."""
        entry = _mock_entry(node_id=node_id)
        other_tenant = uuid.uuid4()
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id, tenant_id=other_tenant),
            ),
        ):
            resp = await client.patch(
                f"/api/v1/documents/{entry.id}",
                json={"material_role": "methodological"},
            )
        assert resp.status_code == 404


class TestIntakeAdmissionWiring:
    """DD-3.3c-I-B route wiring: hint on success, video duration probe, 422."""

    async def test_processing_estimate_attached_on_success(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """A successful create carries the queue-wait hint from settings."""
        from course_supporter.config import get_settings

        s = get_settings()
        entry = _mock_entry(node_id=node_id)
        job = _mock_job()
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "text",
                    "source_url": "https://example.com/doc.md",
                },
            )
        assert resp.status_code == 201
        estimate = resp.json()["processing_estimate"]
        assert estimate["max_hours"] == s.intake_hint_max_hours
        assert estimate["check_after_hours"] == s.intake_hint_check_after_hours

    async def test_video_url_probe_failure_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """An unresolvable video URL yields a clear 422, not a 500."""
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(
                PROBE_URL_FUNC,
                new_callable=AsyncMock,
                side_effect=UnsupportedFormatError("video unavailable"),
            ),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "video",
                    "source_url": "https://youtu.be/private",
                },
            )
        assert resp.status_code == 422
        assert resp.json()["detail"]["code"] == "INTAKE_DURATION_UNAVAILABLE"

    async def test_video_url_probe_processing_error_returns_503(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """A transient probe failure (timeout/missing binary) → 503, not 500."""
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch(
                PROBE_URL_FUNC,
                new_callable=AsyncMock,
                side_effect=ProcessingError("yt-dlp metadata timed out"),
            ),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "video",
                    "source_url": "https://youtu.be/slow",
                },
            )
        assert resp.status_code == 503
        assert resp.json()["detail"]["code"] == "INTAKE_PROBE_UNAVAILABLE"
        assert "retry-after" in {k.lower() for k in resp.headers}

    async def test_video_url_probed_duration_threaded_to_enqueue(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        """The probed video duration reaches enqueue_ingestion."""
        entry = _mock_entry(node_id=node_id, source_type="video")
        job = _mock_job()
        enqueue_mock = AsyncMock(return_value=job)
        with (
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(AuthoredDocumentRepository, "create", return_value=entry),
            patch(PROBE_URL_FUNC, new_callable=AsyncMock, return_value=7200.0),
            patch(ENQUEUE_FUNC, enqueue_mock),
        ):
            resp = await client.post(
                f"/api/v1/nodes/{node_id}/documents",
                data={
                    "source_type": "video",
                    "source_url": "https://youtu.be/lecture",
                },
            )
        assert resp.status_code == 201
        assert enqueue_mock.call_args.kwargs["duration_sec"] == 7200.0


class TestIntakeVideoDurationErrorMapping:
    """`_intake_video_duration_sec` maps probe failures to clean statuses."""

    async def test_non_video_skips_probe_returns_none(self) -> None:
        # No probe patched — a probe call would raise; None proves it's skipped.
        assert await _intake_video_duration_sec(SourceType.TEXT, url="x") is None

    async def test_video_without_content_or_url_raises(self) -> None:
        """Defensive invariant: a video probe with neither input is a hard error."""
        with pytest.raises(ValueError, match="content or url"):
            await _intake_video_duration_sec(SourceType.VIDEO)

    async def test_url_processing_error_maps_to_503(self) -> None:
        with (
            patch(
                PROBE_URL_FUNC,
                new_callable=AsyncMock,
                side_effect=ProcessingError("yt-dlp timed out"),
            ),
            pytest.raises(HTTPException) as exc_info,
        ):
            await _intake_video_duration_sec(SourceType.VIDEO, url="https://youtu.be/x")
        exc = exc_info.value
        assert exc.status_code == 503
        assert isinstance(exc.detail, dict)
        assert exc.detail["code"] == "INTAKE_PROBE_UNAVAILABLE"
        assert exc.headers is not None and "Retry-After" in exc.headers

    async def test_bytes_processing_error_maps_to_503(self) -> None:
        with (
            patch(
                PROBE_BYTES_FUNC,
                new_callable=AsyncMock,
                side_effect=ProcessingError("ffprobe binary not found"),
            ),
            pytest.raises(HTTPException) as exc_info,
        ):
            await _intake_video_duration_sec(
                SourceType.VIDEO, content=b"fake", filename="a.mp4"
            )
        assert exc_info.value.status_code == 503
        assert exc_info.value.detail["code"] == "INTAKE_PROBE_UNAVAILABLE"

    async def test_unsupported_format_still_maps_to_422(self) -> None:
        """The input-error branch is unchanged (subclass matched first)."""
        with (
            patch(
                PROBE_URL_FUNC,
                new_callable=AsyncMock,
                side_effect=UnsupportedFormatError("private video"),
            ),
            pytest.raises(HTTPException) as exc_info,
        ):
            await _intake_video_duration_sec(
                SourceType.VIDEO, url="https://youtu.be/private"
            )
        assert exc_info.value.status_code == 422
        assert exc_info.value.detail["code"] == "INTAKE_DURATION_UNAVAILABLE"


_CONFIRM_PROPOSAL = {
    "files": {
        "src/app.ts": {"role": "full", "reason": "custom_source"},
        "package-lock.json": {"role": "structure_only", "reason": "lockfile"},
    },
    "tree_digest": "digest-1",
    "computed_at": "t0",
}
_CONFIRM_BODY = {
    "files": {"src/app.ts": "full", "package-lock.json": "structure_only"},
    "tree_digest": "digest-1",
}


class TestConfirmFileRoles:
    """POST /api/v1/documents/{did}/file-roles (№21 confirm-api)."""

    def _entry(self, node_id: uuid.UUID) -> MagicMock:
        entry = _mock_entry(node_id=node_id, source_type="code", state="ready")
        entry.file_roles = {"proposal": _CONFIRM_PROPOSAL}
        return entry

    async def test_happy_path_writes_decision_and_enqueues(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        entry = self._entry(node_id)
        job = _mock_job()
        store = AsyncMock(return_value=entry)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(
                AuthoredDocumentRepository, "store_file_roles_decision", store
            ),
            patch.object(
                JobCancellationService,
                "cancel_jobs_for_entities",
                new_callable=AsyncMock,
            ),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=job),
        ):
            resp = await client.post(
                f"/api/v1/documents/{entry.id}/file-roles", json=_CONFIRM_BODY
            )
        assert resp.status_code == 200
        assert resp.json()["job_id"] == str(job.id)
        # Decision written with the full file set + confirmed digest + a stamp.
        decision = store.call_args.kwargs["decision"]
        assert set(decision["files"]) == {"src/app.ts", "package-lock.json"}
        assert decision["tree_digest"] == "digest-1"
        assert "decided_at" in decision

    async def test_auxiliary_role_is_legal(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        entry = self._entry(node_id)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
            patch.object(
                AuthoredDocumentRepository,
                "store_file_roles_decision",
                new_callable=AsyncMock,
                return_value=entry,
            ),
            patch.object(
                JobCancellationService,
                "cancel_jobs_for_entities",
                new_callable=AsyncMock,
            ),
            patch(ENQUEUE_FUNC, new_callable=AsyncMock, return_value=_mock_job()),
        ):
            resp = await client.post(
                f"/api/v1/documents/{entry.id}/file-roles",
                json={
                    "files": {
                        "src/app.ts": "auxiliary",
                        "package-lock.json": "structure_only",
                    },
                    "tree_digest": "digest-1",
                },
            )
        assert resp.status_code == 200

    async def test_digest_mismatch_returns_409_stale(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        entry = self._entry(node_id)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.post(
                f"/api/v1/documents/{entry.id}/file-roles",
                json={**_CONFIRM_BODY, "tree_digest": "STALE"},
            )
        assert resp.status_code == 409
        assert resp.json()["detail"]["code"] == "file_roles_stale"

    async def test_extra_path_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        entry = self._entry(node_id)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.post(
                f"/api/v1/documents/{entry.id}/file-roles",
                json={
                    "files": {**_CONFIRM_BODY["files"], "extra.ts": "full"},
                    "tree_digest": "digest-1",
                },
            )
        assert resp.status_code == 422
        assert resp.json()["detail"]["code"] == "file_roles_incomplete"

    async def test_partial_coverage_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        entry = self._entry(node_id)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.post(
                f"/api/v1/documents/{entry.id}/file-roles",
                json={"files": {"src/app.ts": "full"}, "tree_digest": "digest-1"},
            )
        assert resp.status_code == 422
        assert resp.json()["detail"]["code"] == "file_roles_incomplete"

    async def test_invalid_role_token_returns_422(
        self, client: AsyncClient, node_id: uuid.UUID
    ) -> None:
        entry = self._entry(node_id)
        with (
            patch.object(AuthoredDocumentRepository, "get_by_id", return_value=entry),
            patch.object(
                CourseNodeRepository,
                "get_by_id",
                return_value=_mock_node(node_id=node_id),
            ),
        ):
            resp = await client.post(
                f"/api/v1/documents/{entry.id}/file-roles",
                json={
                    "files": {
                        "src/app.ts": "bogus",
                        "package-lock.json": "structure_only",
                    },
                    "tree_digest": "digest-1",
                },
            )
        assert resp.status_code == 422  # Pydantic Literal rejection

    async def test_not_found_returns_404(self, client: AsyncClient) -> None:
        with patch.object(AuthoredDocumentRepository, "get_by_id", return_value=None):
            resp = await client.post(
                f"/api/v1/documents/{uuid.uuid4()}/file-roles", json=_CONFIRM_BODY
            )
        assert resp.status_code == 404
